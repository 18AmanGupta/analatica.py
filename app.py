import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from sklearn.model_selection import train_test_split, GridSearchCV
from sklearn.ensemble import RandomForestRegressor, RandomForestClassifier
from sklearn.metrics import mean_absolute_error, r2_score, accuracy_score, classification_report, precision_score, recall_score, f1_score
from statsmodels.tsa.arima.model import ARIMA
from fpdf import FPDF
import io
import shap
import tensorflow as tf
from tensorflow import keras
from pptx import Presentation
from pptx.util import Inches
import xlsxwriter

# Function to Load Data
def load_data(uploaded_file):
    df = pd.read_csv(uploaded_file) if uploaded_file.name.endswith('.csv') else pd.read_excel(uploaded_file)
    return df

# Function for Data Cleaning
def clean_data(df):
    df.fillna(method='ffill', inplace=True)
    df.fillna(method='bfill', inplace=True)
    df.fillna(df.select_dtypes(include=[np.number]).median(), inplace=True)
    
    for col in df.columns:
        if df[col].dtype == 'object':
            try:
                df[col] = pd.to_datetime(df[col])
                df["Year"] = df[col].dt.year
                df["Month"] = df[col].dt.month
                df["Day"] = df[col].dt.day
                df["Hour"] = df[col].dt.hour
                df.drop(columns=[col], inplace=True)
            except:
                continue
    
    return df

# Function for Statistical Insights
def statistical_insights(df):
    st.subheader("📊 Statistical Insights")
    st.write(df.describe())

# Function for Custom Dashboard
def custom_dashboard(df):
    st.subheader("📊 Custom Dashboard")
    selected_columns = st.multiselect("Select Columns for Dashboard", df.columns)
    if selected_columns:
        st.line_chart(df[selected_columns])

# Function to Determine Task Type
def determine_task(df):
    target_col = st.selectbox("Select Target Column", df.columns)
    if df[target_col].dtype == 'object':
        return "classification", target_col
    else:
        return "regression", target_col

# Function to Train ML Model with Hyperparameter Tuning and SHAP Explanation
def train_ml_model(df, task_type, target_col):
    X = df.drop(columns=[target_col])
    y = df[target_col]
    X = pd.get_dummies(X, drop_first=True)
    X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
    
    if task_type == "regression":
        param_grid = {"n_estimators": [50, 100, 200], "max_depth": [None, 10, 20]}
        model = GridSearchCV(RandomForestRegressor(random_state=42), param_grid, cv=3)
        model.fit(X_train, y_train)
        y_pred = model.best_estimator_.predict(X_test)
        mae = mean_absolute_error(y_test, y_pred)
        r2 = r2_score(y_test, y_pred)
        explainer = shap.Explainer(model.best_estimator_, X_train)
        shap_values = explainer(X_test)
        return model.best_estimator_, f"Mean Absolute Error: {mae:.2f}", f"R² Score: {r2:.2f}", shap_values
    else:
        param_grid = {"n_estimators": [50, 100, 200], "max_depth": [None, 10, 20]}
        model = GridSearchCV(RandomForestClassifier(random_state=42), param_grid, cv=3)
        model.fit(X_train, y_train)
        y_pred = model.best_estimator_.predict(X_test)
        accuracy = accuracy_score(y_test, y_pred)
        precision = precision_score(y_test, y_pred, average='weighted', zero_division=0)
        recall = recall_score(y_test, y_pred, average='weighted', zero_division=0)
        f1 = f1_score(y_test, y_pred, average='weighted', zero_division=0)
        report = classification_report(y_test, y_pred)
        explainer = shap.Explainer(model.best_estimator_, X_train)
        shap_values = explainer(X_test)
        return model.best_estimator_, f"Accuracy: {accuracy:.2f}", f"Precision: {precision:.2f}", f"Recall: {recall:.2f}", f"F1 Score: {f1:.2f}", report, shap_values

# Function to Export Reports with Charts
def export_report(df, metrics, format):
    if format == "PDF":
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.cell(200, 10, txt="Data Analysis Report", ln=True, align='C')
        pdf.ln(10)
        pdf.multi_cell(0, 10, txt=str(df.describe()))
        pdf.ln(10)
        pdf.multi_cell(0, 10, txt=str(metrics))
        pdf.output("report.pdf")
        st.success("Report exported as PDF")
    elif format == "Excel":
        excel_buffer = io.BytesIO()
        with pd.ExcelWriter(excel_buffer, engine='xlsxwriter') as writer:
            df.to_excel(writer, sheet_name='Data', index=False)
            pd.DataFrame(metrics).to_excel(writer, sheet_name='Metrics', index=False)
        st.download_button(label="Download Excel Report", data=excel_buffer.getvalue(), file_name="report.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    elif format == "PowerPoint":
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Data Analysis Report"
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        text_frame = text_box.text_frame
        text_frame.text = "Statistical Insights:\n" + str(df.describe())
        ppt.save("report.pptx")
        st.success("Report exported as PowerPoint (report.pptx)")

# Streamlit UI
def main():
    st.title("📊 General-Purpose Data Analysis & Prediction App")
    uploaded_file = st.file_uploader("Upload a CSV or Excel file", type=["csv", "xlsx"])
    
    if uploaded_file:
        df = load_data(uploaded_file)
        df = clean_data(df)
        
        st.subheader("💒 Preview of Data")
        num_rows = st.slider("Select number of rows to display", min_value=5, max_value=len(df), value=10)
        st.write(df.head(num_rows))
        
        statistical_insights(df)
        custom_dashboard(df)
        
        task_type, target_col = determine_task(df)
        model, *metrics, shap_values = train_ml_model(df, task_type, target_col)
        
        st.subheader("⚡ Machine Learning Prediction")
        for metric in metrics:
            st.write(metric)
        
        format_option = st.selectbox("Select Report Format", ["PDF", "Excel", "PowerPoint"])
        if st.button("Export Report"):
            export_report(df, metrics, format_option)

if __name__ == "__main__":
    main()









